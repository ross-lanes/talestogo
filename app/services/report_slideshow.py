"""
PowerPoint Slideshow Export Service
Generates PowerPoint presentations with:
- Dashboard image (Key Metrics)
- All sub-analysis charts from Analytics pages
- Recommendations as formatted text slides
"""

import os
import io
import re
from typing import Optional, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from sqlalchemy.orm import Session


def generate_slideshow(
    markdown_content: str,
    title: str,
    db: Session,
    user_id: int,
    brand_id: Optional[int] = None
) -> io.BytesIO:
    """
    Generate a PowerPoint slideshow with dashboard, analytics charts, and recommendations.

    Args:
        markdown_content: The markdown report content
        title: Report title
        db: Database session
        user_id: User ID
        brand_id: Brand ID

    Returns:
        BytesIO object containing the PowerPoint file
    """
    from app import models
    from app.services.chart_generator import ensure_charts_directory

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get brand name
    brand_name = "Your Brand"
    if brand_id:
        brand = db.query(models.BrandInfo).filter(models.BrandInfo.id == brand_id).first()
        if brand:
            brand_name = brand.brand_name

    # RobotRachel brand colors
    BRAND_BLUE = RGBColor(128, 161, 212)    # #80a1d4 - Primary blue
    BRAND_PURPLE = RGBColor(102, 87, 117)   # #665775 - Secondary purple
    BRAND_TEAL = RGBColor(117, 201, 200)    # #75c9c8 - Accent teal

    # === Slide 1: Title Slide ===
    _create_title_slide(prs, title, brand_name, BRAND_BLUE)

    # === Slide 2: Executive Summary ===
    exec_summary = _extract_executive_summary(markdown_content)
    if exec_summary:
        _create_text_slide(prs, "Executive Summary", exec_summary, BRAND_BLUE)

    # === Slide 3: Dashboard Image (Key Metrics) ===
    dashboard_path = _find_dashboard_image(user_id, brand_id, brand_name)
    if dashboard_path and os.path.exists(dashboard_path):
        _create_image_slide(prs, "Key Metrics Dashboard", dashboard_path, BRAND_BLUE)

    # === Analytics Charts Slides ===
    analytics_charts = _find_all_analytics_charts(brand_name)
    for chart_title, chart_path in analytics_charts:
        if os.path.exists(chart_path):
            _create_image_slide(prs, chart_title, chart_path, BRAND_BLUE)

    # === Recommendations Slides ===
    recommendations = _extract_recommendations(markdown_content)
    for rec_title, rec_bullets in recommendations:
        _create_recommendations_slide(prs, rec_title, rec_bullets, BRAND_BLUE)

    # === Final Slide: Questions ===
    _create_final_slide(prs, BRAND_TEAL)

    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output


def _create_title_slide(prs: Presentation, title: str, brand_name: str, color: RGBColor):
    """Create the title slide with blue background."""
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)

    # Add blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"AI Reputation Analysis for {brand_name}"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER


def _create_image_slide(prs: Presentation, title: str, image_path: str, color: RGBColor):
    """
    Create a slide with a title and properly sized/centered image.
    Maintains aspect ratio and ensures image fits within available space.
    """
    from PIL import Image

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color

    # Add image with proper sizing
    try:
        # Get image dimensions
        img = Image.open(image_path)
        img_width, img_height = img.size
        img_aspect = img_width / img_height

        # Define available space on slide
        # Slide is 10" wide x 7.5" tall
        # Title takes up 0.3" + 0.7" = 1" from top
        # Leave 0.5" margins on sides and 0.5" at bottom
        available_width = Inches(9)   # 10 - 0.5 (left) - 0.5 (right)
        available_height = Inches(6)  # 7.5 - 1 (title) - 0.5 (bottom margin)
        available_top = Inches(1.2)
        available_left = Inches(0.5)

        # Calculate dimensions that fit within available space while maintaining aspect ratio
        available_aspect = available_width / available_height

        if img_aspect > available_aspect:
            # Image is wider than available space - constrain by width
            final_width = available_width
            final_height = available_width / img_aspect
        else:
            # Image is taller than available space - constrain by height
            final_height = available_height
            final_width = available_height * img_aspect

        # Center the image horizontally
        left = available_left + (available_width - final_width) / 2

        # Add the image
        slide.shapes.add_picture(
            image_path,
            left,
            available_top,
            width=final_width,
            height=final_height
        )

    except Exception as e:
        print(f"Error adding image {image_path}: {e}")


def _create_text_slide(prs: Presentation, title: str, content: str, color: RGBColor):
    """Create a slide with title and bulleted text content."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color

    # Add bulleted content
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    # Remove markdown formatting
    clean_content = re.sub(r'\*\*(.*?)\*\*', r'\1', content)

    # Split into sentences and create bullets
    sentences = [s.strip() + '.' for s in clean_content.split('.') if s.strip() and len(s.strip()) > 20]

    # Take first 6 key points
    key_points = sentences[:6]

    for i, point in enumerate(key_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = f"• {point.strip()}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0


def _create_recommendations_slide(prs: Presentation, title: str, bullets: List[str], color: RGBColor):
    """Create a slide with recommendation title and bulleted actions."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = color

    # Add bulleted recommendations
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = f"• {bullet.strip()}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0


def _create_final_slide(prs: Presentation, color: RGBColor):
    """Create the final 'Questions?' slide."""
    final_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(final_slide_layout)

    # Add teal background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

    # Add thank you text
    thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Questions?"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(48)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_you_para.alignment = PP_ALIGN.CENTER


def _extract_executive_summary(markdown_content: str) -> Optional[str]:
    """Extract executive summary from markdown content."""
    exec_summary_match = re.search(
        r'## Executive Summary\n\n(.*?)(?=\n\n!|\n\n##|\Z)',
        markdown_content,
        re.DOTALL
    )
    if exec_summary_match:
        return exec_summary_match.group(1).strip()
    return None


def _extract_recommendations(markdown_content: str) -> List[Tuple[str, List[str]]]:
    """
    Extract recommendations section and parse into individual recommendation slides.
    Returns list of (title, bullet_points) tuples.
    """
    recommendations = []

    # Find the Recommendations section
    rec_match = re.search(
        r'### 7\. Recommendations\n\n(.*?)(?=\n\n##|\Z)',
        markdown_content,
        re.DOTALL
    )

    if not rec_match:
        return recommendations

    rec_content = rec_match.group(1).strip()

    # Split by paragraphs to get individual recommendations
    # Each recommendation typically starts with a number or bullet
    paragraphs = [p.strip() for p in rec_content.split('\n\n') if p.strip()]

    current_rec_bullets = []
    for para in paragraphs:
        # Check if this is a numbered recommendation (starts with digit)
        if re.match(r'^\d+\.', para):
            # Save previous recommendation if it exists
            if current_rec_bullets:
                # Take first line as title, rest as bullets
                first_line = current_rec_bullets[0]
                # Extract title from "1. **Title**: description" format
                title_match = re.match(r'^\d+\.\s*\*\*(.*?)\*\*:', first_line)
                if title_match:
                    title = title_match.group(1)
                    # Get rest of first line as first bullet
                    rest_of_line = re.sub(r'^\d+\.\s*\*\*.*?\*\*:\s*', '', first_line)
                    bullets = [rest_of_line] + current_rec_bullets[1:]
                    recommendations.append((title, bullets[:6]))  # Limit to 6 bullets
                else:
                    # Fallback: use first line as title
                    title = re.sub(r'^\d+\.\s*', '', first_line)
                    recommendations.append((title, current_rec_bullets[1:][:6]))

            # Start new recommendation
            current_rec_bullets = [para]
        else:
            # Add to current recommendation
            if current_rec_bullets:
                # Split by lines starting with "- " or "* "
                lines = para.split('\n')
                for line in lines:
                    if line.strip().startswith(('- ', '* ', '• ')):
                        clean_line = re.sub(r'^[-*•]\s*', '', line.strip())
                        current_rec_bullets.append(clean_line)

    # Don't forget the last recommendation
    if current_rec_bullets:
        first_line = current_rec_bullets[0]
        title_match = re.match(r'^\d+\.\s*\*\*(.*?)\*\*:', first_line)
        if title_match:
            title = title_match.group(1)
            rest_of_line = re.sub(r'^\d+\.\s*\*\*.*?\*\*:\s*', '', first_line)
            bullets = [rest_of_line] + current_rec_bullets[1:]
            recommendations.append((title, bullets[:6]))
        else:
            title = re.sub(r'^\d+\.\s*', '', first_line)
            recommendations.append((title, current_rec_bullets[1:][:6]))

    return recommendations


def _find_dashboard_image(user_id: int, brand_id: Optional[int], brand_name: str) -> Optional[str]:
    """
    Find the Dashboard/Key Metrics image from the report_charts directory.
    The Dashboard download feature saves files as: KeyMetrics_{BrandName}_{date}.png

    NOTE: Dashboard images are currently saved locally by users, not uploaded to server.
    This function will return None until we add Dashboard upload functionality.
    """
    import glob

    charts_dir = os.path.join('frontend', 'public', 'report_charts')
    if not os.path.exists(charts_dir):
        return None

    # Format brand name for filename matching (spaces to underscores, no special chars)
    brand_name_formatted = brand_name.replace(' ', '_').replace('-', '_')

    # Look for KeyMetrics files for this brand
    pattern = os.path.join(charts_dir, f"KeyMetrics_{brand_name_formatted}_*.png")
    files = glob.glob(pattern)

    if files:
        # Return the most recent file
        return max(files, key=os.path.getmtime)

    return None


def _find_all_analytics_charts(brand_name: str) -> List[Tuple[str, str]]:
    """
    Find all analytics charts from the report_charts directory for a specific brand.
    These are the actual website screenshots saved when users click "Download as Image".
    Returns list of (title, path) tuples in the correct order.

    Args:
        brand_name: Brand name to match in filenames (e.g., "Princeton Engineering")

    Returns:
        List of (chart_title, chart_path) tuples
    """
    import glob
    import re

    charts = []
    found_charts = {}

    # Chart titles mapping (in display order)
    chart_order = [
        ('mention_rate', 'Brand Mention Rate'),
        ('brand_mentions_trend', 'Brand Mentions Trend'),
        ('platform_comparison', 'Platform Comparison'),
        ('positioning', 'Brand Positioning'),
        ('positioning_trend', 'Positioning Trend'),
        ('sentiment', 'Sentiment Distribution'),
        ('sentiment_trend', 'Sentiment Trend'),
        ('share_of_voice', 'Share of Voice'),
        ('share_of_voice_trend', 'Share of Voice Trend'),
        ('descriptor_performance', 'Top Descriptors'),
        ('competitor_threats', 'Competitive Threats'),
    ]

    # Search for all chart files in the report_charts directory
    charts_dir = os.path.join('frontend', 'public', 'report_charts')
    if not os.path.exists(charts_dir):
        return charts

    # Format brand name for filename matching (spaces to underscores)
    brand_name_formatted = brand_name.replace(' ', '_').replace('-', '_')

    # Get all PNG files for this brand
    all_brand_files = glob.glob(os.path.join(charts_dir, f"{brand_name_formatted}_*.png"))

    # Group files by chart type and find the most recent for each
    chart_files_by_type = {}
    for chart_path in all_brand_files:
        filename = os.path.basename(chart_path)

        # Match to chart types
        for key, title in chart_order:
            if key in filename.lower():
                # Extract timestamp from filename (last part before .png)
                # Format: BrandName_charttype_YYYYMMDD_HHMMSS.png
                timestamp_match = re.search(r'_(\d{8}_\d{6})\.png$', filename)
                if timestamp_match:
                    timestamp = timestamp_match.group(1)

                    # Keep track of all files for this chart type
                    if key not in chart_files_by_type:
                        chart_files_by_type[key] = []
                    chart_files_by_type[key].append((timestamp, chart_path, title))
                break

    # For each chart type, select the most recent file
    for key, files in chart_files_by_type.items():
        # Sort by timestamp (most recent first)
        files.sort(key=lambda x: x[0], reverse=True)
        most_recent = files[0]
        found_charts[key] = (most_recent[2], most_recent[1])  # (title, path)

    # Return charts in the specified order
    for key, _ in chart_order:
        if key in found_charts:
            charts.append(found_charts[key])

    return charts
