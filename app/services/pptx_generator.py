"""
PowerPoint deck generation service for personas
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import List, Dict, Any, Optional
import os


class PersonaDeckGenerator:
    """Generate PowerPoint presentations for persona decks"""

    def __init__(self):
        """Initialize the generator"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def add_title_slide(
        self, brand_name: str, subtitle: str = "Healthcare Professional & Patient Personas",
        primary_color: str = "#1976d2", logo_url: Optional[str] = None
    ):
        """Add a title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Brand name
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = brand_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        rgb = self.hex_to_rgb(primary_color)
        title_para.font.color.rgb = RGBColor(*rgb)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Add logo if provided
        if logo_url and os.path.exists(logo_url):
            try:
                slide.shapes.add_picture(logo_url, Inches(4), Inches(0.5), height=Inches(1))
            except:
                pass  # If logo fails, continue without it

    def add_patient_persona_slide(
        self, persona: Dict[str, Any], brand_name: str,
        primary_color: str = "#1976d2", secondary_color: str = "#dc004e"
    ):
        """Add a patient persona slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Colors
        primary_rgb = self.hex_to_rgb(primary_color)
        secondary_rgb = self.hex_to_rgb(secondary_color)

        # Left sidebar background
        left_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(3), Inches(7.5)
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(240, 240, 240)
        left_bg.line.color.rgb = RGBColor(200, 200, 200)

        # Name (top of slide)
        name_box = slide.shapes.add_textbox(Inches(3.2), Inches(0.3), Inches(6.5), Inches(0.6))
        name_frame = name_box.text_frame
        name_frame.text = persona.get('name', 'Unknown')
        name_para = name_frame.paragraphs[0]
        name_para.font.size = Pt(36)
        name_para.font.bold = True
        name_para.font.color.rgb = RGBColor(*primary_rgb)

        # Quote (left sidebar)
        quote_box = slide.shapes.add_textbox(Inches(0.3), Inches(2), Inches(2.4), Inches(1.2))
        quote_frame = quote_box.text_frame
        quote_frame.text = f'"{persona.get("quote", "")}"'
        quote_para = quote_frame.paragraphs[0]
        quote_para.font.size = Pt(12)
        quote_para.font.italic = True
        quote_para.font.color.rgb = RGBColor(80, 80, 80)
        quote_para.alignment = PP_ALIGN.CENTER
        quote_frame.word_wrap = True

        # Key demographics (left sidebar)
        demo_data = [
            ("AGE", str(persona.get('age', 'N/A'))),
            ("LOCATION", persona.get('location', 'N/A')),
            ("STATUS", persona.get('family_status', 'N/A')),
            ("OCCUPATION", persona.get('occupation', 'N/A'))
        ]

        y_pos = 3.5
        for label, value in demo_data:
            # Label
            label_box = slide.shapes.add_textbox(Inches(0.3), Inches(y_pos), Inches(2.4), Inches(0.25))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(10)
            label_para.font.bold = True
            label_para.font.color.rgb = RGBColor(60, 60, 60)

            # Value
            value_box = slide.shapes.add_textbox(Inches(0.3), Inches(y_pos + 0.25), Inches(2.4), Inches(0.35))
            value_frame = value_box.text_frame
            value_frame.text = str(value)[:40]  # Truncate if too long
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(11)
            value_para.font.color.rgb = RGBColor(40, 40, 40)
            value_frame.word_wrap = True

            y_pos += 0.65

        # Personality traits (left sidebar)
        traits = persona.get('personality_traits', '').split(',')
        if traits and traits[0]:
            y_pos = 6.2
            for i, trait in enumerate(traits[:4]):
                trait_box = slide.shapes.add_shape(
                    1, Inches(0.3), Inches(y_pos), Inches(1.1), Inches(0.3)
                )
                trait_box.fill.solid()
                trait_box.fill.fore_color.rgb = RGBColor(*secondary_rgb)
                trait_box.line.fill.background()

                trait_text = trait_box.text_frame
                trait_text.text = trait.strip().upper()[:15]
                trait_para = trait_text.paragraphs[0]
                trait_para.font.size = Pt(8)
                trait_para.font.bold = True
                trait_para.font.color.rgb = RGBColor(255, 255, 255)
                trait_para.alignment = PP_ALIGN.CENTER
                trait_text.vertical_anchor = MSO_ANCHOR.MIDDLE

                y_pos += 0.35

        # Main content (right side)
        sections = [
            ("DEMOGRAPHICS", self._format_demographics(persona)),
            ("CLINICAL PROFILE", self._format_clinical_profile(persona)),
            ("GOALS & FEARS", self._format_goals_fears(persona)),
            ("INFORMATION JOURNEY", self._format_info_journey(persona)),
            (f"MARKETING CUES FOR {brand_name.upper()}", self._format_marketing_cues(persona))
        ]

        x_left = 3.3
        x_right = 6.7
        y_pos = 1.1
        row_height = 1.15

        for i, (title, content) in enumerate(sections):
            if i < 3:  # First 3 in left column
                x = x_left
            else:  # Last 2 in right column
                x = x_right
                if i == 3:
                    y_pos = 1.1  # Reset y position for right column

            # Section title
            title_box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(3.2), Inches(0.25))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(10)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(*primary_rgb)

            # Section content
            content_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_pos + 0.28), Inches(3.2), Inches(row_height - 0.28)
            )
            content_frame = content_box.text_frame
            content_frame.text = content
            content_frame.word_wrap = True
            for para in content_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.color.rgb = RGBColor(40, 40, 40)
                para.space_after = Pt(4)

            if i < 3:
                y_pos += row_height

    def add_hcp_persona_slide(
        self, persona: Dict[str, Any], brand_name: str,
        primary_color: str = "#1976d2", secondary_color: str = "#dc004e"
    ):
        """Add an HCP persona slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Colors
        primary_rgb = self.hex_to_rgb(primary_color)
        secondary_rgb = self.hex_to_rgb(secondary_color)

        # Left sidebar background
        left_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(3), Inches(7.5)
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(240, 240, 240)
        left_bg.line.color.rgb = RGBColor(200, 200, 200)

        # Name (top of slide)
        name_box = slide.shapes.add_textbox(Inches(3.2), Inches(0.3), Inches(6.5), Inches(0.6))
        name_frame = name_box.text_frame
        name_frame.text = persona.get('name', 'Unknown')
        name_para = name_frame.paragraphs[0]
        name_para.font.size = Pt(36)
        name_para.font.bold = True
        name_para.font.color.rgb = RGBColor(*primary_rgb)

        # Quote (left sidebar)
        quote_box = slide.shapes.add_textbox(Inches(0.3), Inches(2), Inches(2.4), Inches(1.2))
        quote_frame = quote_box.text_frame
        quote_frame.text = f'"{persona.get("quote", "")}"'
        quote_para = quote_frame.paragraphs[0]
        quote_para.font.size = Pt(11)
        quote_para.font.italic = True
        quote_para.font.color.rgb = RGBColor(80, 80, 80)
        quote_para.alignment = PP_ALIGN.CENTER
        quote_frame.word_wrap = True

        # Key demographics (left sidebar)
        demo_data = [
            ("AGE", str(persona.get('age', 'N/A'))),
            ("JOB TITLE", persona.get('job_title', 'N/A')[:30]),
            ("PRACTICE", persona.get('practice_type', 'N/A')[:30]),
            ("LOCATION", persona.get('location', 'N/A'))
        ]

        y_pos = 3.5
        for label, value in demo_data:
            # Label
            label_box = slide.shapes.add_textbox(Inches(0.3), Inches(y_pos), Inches(2.4), Inches(0.25))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(10)
            label_para.font.bold = True
            label_para.font.color.rgb = RGBColor(60, 60, 60)

            # Value
            value_box = slide.shapes.add_textbox(Inches(0.3), Inches(y_pos + 0.25), Inches(2.4), Inches(0.35))
            value_frame = value_box.text_frame
            value_frame.text = str(value)
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(10)
            value_para.font.color.rgb = RGBColor(40, 40, 40)
            value_frame.word_wrap = True

            y_pos += 0.65

        # Personality traits (left sidebar)
        traits = persona.get('personality_traits', '').split(',')
        if traits and traits[0]:
            y_pos = 6.2
            for i, trait in enumerate(traits[:4]):
                trait_box = slide.shapes.add_shape(
                    1, Inches(0.3), Inches(y_pos), Inches(1.1), Inches(0.3)
                )
                trait_box.fill.solid()
                trait_box.fill.fore_color.rgb = RGBColor(*secondary_rgb)
                trait_box.line.fill.background()

                trait_text = trait_box.text_frame
                trait_text.text = trait.strip().upper()[:15]
                trait_para = trait_text.paragraphs[0]
                trait_para.font.size = Pt(8)
                trait_para.font.bold = True
                trait_para.font.color.rgb = RGBColor(255, 255, 255)
                trait_para.alignment = PP_ALIGN.CENTER
                trait_text.vertical_anchor = MSO_ANCHOR.MIDDLE

                y_pos += 0.35

        # Main content (right side)
        sections = [
            ("ABOUT", persona.get('about', 'N/A')[:400]),
            ("GOALS & MOTIVATIONS", persona.get('motivations', 'N/A')),
            ("PAIN POINTS & FRUSTRATIONS", persona.get('frustrations', 'N/A')),
            (f"HOW THEY VIEW {brand_name.upper()}", persona.get('how_they_view_brand', 'N/A')[:300]),
            ("MARKETING STRATEGY", self._format_hcp_marketing(persona))
        ]

        x_left = 3.3
        x_right = 6.7
        y_pos = 1.1
        row_height = 1.15

        for i, (title, content) in enumerate(sections):
            if i < 3:
                x = x_left
            else:
                x = x_right
                if i == 3:
                    y_pos = 1.1

            # Section title
            title_box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(3.2), Inches(0.25))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(10)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(*primary_rgb)

            # Section content
            content_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_pos + 0.28), Inches(3.2), Inches(row_height - 0.28)
            )
            content_frame = content_box.text_frame
            content_frame.text = content
            content_frame.word_wrap = True
            for para in content_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.color.rgb = RGBColor(40, 40, 40)
                para.space_after = Pt(4)

            if i < 3:
                y_pos += row_height

    def _format_demographics(self, persona: Dict[str, Any]) -> str:
        """Format demographics section for patient"""
        return f"""Age: {persona.get('age', 'N/A')}
Location: {persona.get('location', 'N/A')}
Family Status: {persona.get('family_status', 'N/A')}
Occupation: {persona.get('occupation', 'N/A')}
Tech Savviness: {persona.get('tech_savviness', 'N/A')}"""

    def _format_clinical_profile(self, persona: Dict[str, Any]) -> str:
        """Format clinical profile section"""
        return f"""Scenario: {persona.get('clinical_scenario', 'N/A')[:150]}
Recent Diagnosis: {persona.get('recent_diagnosis', 'N/A')[:100]}
Mindset: {persona.get('mindset', 'N/A')[:100]}"""

    def _format_goals_fears(self, persona: Dict[str, Any]) -> str:
        """Format goals and fears section"""
        goals = persona.get('goals', '')
        fears = persona.get('fears', '')
        return f"""GOALS:\n{goals[:200] if goals else 'N/A'}

FEARS:\n{fears[:200] if fears else 'N/A'}"""

    def _format_info_journey(self, persona: Dict[str, Any]) -> str:
        """Format information journey section"""
        return f"""Channels: {persona.get('information_channels', 'N/A')[:150]}

Key Question: {persona.get('key_question_for_doctor', 'N/A')[:150]}"""

    def _format_marketing_cues(self, persona: Dict[str, Any]) -> str:
        """Format marketing cues section"""
        return f"""Message: {persona.get('marketing_message', 'N/A')[:150]}

Tone: {persona.get('marketing_tone', 'N/A')[:80]}

Call to Action: {persona.get('call_to_action', 'N/A')[:100]}"""

    def _format_hcp_marketing(self, persona: Dict[str, Any]) -> str:
        """Format HCP marketing strategy section"""
        channels = persona.get('marketing_channels', 'N/A')[:100]
        messaging = persona.get('marketing_messaging', 'N/A')[:150]
        return f"""Channels: {channels}

Messaging: {messaging}"""

    def save(self, filepath: str):
        """Save the presentation to a file"""
        self.prs.save(filepath)
        return filepath
